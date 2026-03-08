import html as pyhtml
import mimetypes
import os
import subprocess
import uuid
import xml.etree.ElementTree as ET
import zipfile
from datetime import datetime
from pathlib import Path
from urllib.parse import quote

from docx import Document as DocxDocument
from fastapi import APIRouter, Depends, Form, HTTPException, Query, Request, UploadFile
from fastapi import File as FastAPIFile
from fastapi.responses import FileResponse, HTMLResponse
from pptx import Presentation
from pydantic import BaseModel
from sqlalchemy import func
from sqlalchemy.orm import Session

from proekt.backend.auth import decode_access_token
from proekt.backend.database import SessionLocal
from proekt.backend.models import Comment, Rating, User
from proekt.backend.models import File as FileModel

# Папка для зберігання файлів
UPLOAD_DIR = Path(__file__).resolve().parent.parent / "uploads"
UPLOAD_DIR.mkdir(exist_ok=True)

router = APIRouter(prefix="/files", tags=["files"])


def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


def _require_token(token: str):
    if not token:
        raise HTTPException(status_code=401, detail="Токен не предоставлен")
    try:
        payload = decode_access_token(token)
    except ValueError:
        raise HTTPException(status_code=401, detail="Неверный или просроченный токен")
    return payload


# Pydantic модель для переименування
class RenameRequest(BaseModel):
    name: str


@router.post("/create-folder")
def create_folder(
    token: str = Query(...),
    folder_name: str = Form(...),
    comment: str = Form(None),  # ✅ ДОДАТИ
    is_public: bool = Form(False),
    db: Session = Depends(get_db),
):
    try:
        payload = _require_token(token)

        user_id = payload.get("user_id")
        user = db.query(User).filter(User.id == user_id).first()

        if not user:
            raise HTTPException(status_code=401, detail="Пользователь не найден")

        if not folder_name or folder_name.strip() == "":
            raise HTTPException(
                status_code=400, detail="Имя папки не может быть пустым"
            )

        folder_id = str(uuid.uuid4())

        # створити фізичну папку під uploads
        folder_path = UPLOAD_DIR / folder_id
        folder_path.mkdir(parents=True, exist_ok=True)

        db_folder = FileModel(
            user_id=user.id,
            folder_id=folder_id,
            filename=folder_name,
            file_path=str(folder_path),  # зберігаємо шлях до папки
            file_type="folder",
            file_size=0,
            is_public=is_public,
            comment=comment,
            created_at=datetime.utcnow(),
        )
        db.add(db_folder)
        db.commit()
        db.refresh(db_folder)

        return {"id": folder_id, "name": folder_name, "type": "folder"}
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


def get_file_type(filename: str) -> str:
    """Визначити тип файлу"""
    ext = Path(filename).suffix.lower()
    video_exts = {".mp4", ".avi", ".mov", ".mkv", ".flv", ".wmv"}
    audio_exts = {".mp3", ".wav", ".aac", ".flac", ".wma", ".m4a"}
    image_exts = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".svg", ".webp"}

    if ext in video_exts:
        return "video"
    elif ext in audio_exts:
        return "audio"
    elif ext in image_exts:
        return "image"
    else:
        return "document"


def convert_ppt_to_pdf(src_path: str, dst_path: str) -> None:
    """Convert PPT/PPTX to PDF.

    Tries:
    1) LibreOffice (soffice)
    2) Microsoft PowerPoint COM automation (Windows)
    3) Fallback: render text from slides into a simple PDF using reportlab
    """

    # 1) Try LibreOffice (if installed)
    try:
        subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(Path(dst_path).parent),
                src_path,
            ],
            capture_output=True,
            check=True,
            text=True,
        )
        if Path(dst_path).exists():
            return
        libre_error = "LibreOffice conversion did not produce output"
    except FileNotFoundError:
        libre_error = "LibreOffice (soffice) not found."
    except subprocess.CalledProcessError as e:
        libre_error = f"LibreOffice conversion failed: {e.stderr.strip()}"

    # 2) Try PowerPoint COM automation (Windows)
    try:
        import win32com.client

        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = False
        pres = ppt_app.Presentations.Open(str(src_path), WithWindow=False)
        pres.SaveAs(str(dst_path), 32)  # 32 = ppSaveAsPDF
        pres.Close()
        ppt_app.Quit()
        if Path(dst_path).exists():
            return
        ppt_error = "PowerPoint conversion did not produce output"
    except Exception as e:
        ppt_error = str(e)

    # 3) Fallback: render slide text into a basic PDF (no images/layout)
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas

        prs = Presentation(src_path)
        c = canvas.Canvas(str(dst_path), pagesize=letter)
        width, height = letter
        line_height = 14
        text_margin = 40

        for slide in prs.slides:
            y = height - text_margin
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    for line in shape.text.splitlines():
                        c.drawString(text_margin, y, line[:120])
                        y -= line_height
                        if y < text_margin:
                            c.showPage()
                            y = height - text_margin
            c.showPage()
        c.save()
        if Path(dst_path).exists():
            return
        raise RuntimeError("Reportlab conversion did not produce output")
    except Exception as e:
        raise RuntimeError(
            f"Cannot convert PPT/PPTX to PDF. LibreOffice: {libre_error}; PowerPoint: {ppt_error}; Text fallback: {e}"
        )


# ============ НОВИЙ ENDPOINT ============
@router.get("/folders")
def get_user_folders(token: str = Query(...), db: Session = Depends(get_db)):
    """Отримати всі папки користувача"""
    try:
        payload = _require_token(token)
        user_id = payload.get("user_id")
        user = db.query(User).filter(User.id == user_id).first()

        if not user:
            raise HTTPException(status_code=401, detail="Пользователь не найден")

        folders = (
            db.query(FileModel)
            .filter(FileModel.user_id == user.id, FileModel.file_type == "folder")
            .all()
        )

        # Підрахунок файлів для кожної папки
        folder_data = []
        for f in folders:
            file_count = (
                db.query(func.count(FileModel.id))
                .filter(
                    FileModel.folder_id == f.folder_id, FileModel.file_type != "folder"
                )
                .scalar()
            )
            folder_data.append(
                {
                    "id": f.folder_id,
                    "name": f.filename,
                    "is_public": f.is_public,
                    "comment": f.comment,
                    "file_count": file_count,
                    "created_at": f.created_at.isoformat() if f.created_at else None,
                }
            )

        return folder_data

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/folder-info/{folder_id}")
def get_folder_info(
    folder_id: str, token: str = Query(None), db: Session = Depends(get_db)
):
    """Отримати інформацію про папку (публічну або приватну з токеном)"""
    try:
        user_id = None

        # Спробуємо отримати user_id з токена, якщо він є
        if token:
            try:
                payload = decode_access_token(token)
                user_id = payload.get("user_id")
            except ValueError:
                pass  # Токен невірний, але можемо показати публічні папки

        # Отримуємо папку
        folder = (
            db.query(FileModel)
            .filter(FileModel.folder_id == folder_id, FileModel.file_type == "folder")
            .first()
        )

        if not folder:
            raise HTTPException(status_code=404, detail="Папка не найдена")

        # Перевіряємо доступ
        is_owner = user_id == folder.user_id
        is_public = folder.is_public

        if not is_owner and not is_public:
            raise HTTPException(status_code=403, detail="Нет доступа к папке")

        user = db.query(User).filter(User.id == folder.user_id).first()
        if not user:
            raise HTTPException(status_code=404, detail="Пользователь не найден")

        # Формуємо ім'я власника
        owner_name = f"{user.name or ''} {user.surname or ''}".strip() or user.email

        return {
            "name": folder.filename,
            "comment": folder.comment or "",
            "owner": owner_name,
            "is_public": is_public,
            "date": folder.created_at.isoformat(),
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# =====================================
@router.put("/rename/{file_id}")
def rename_file(
    file_id: int,
    request: RenameRequest,
    token: str = Query(...),
    db: Session = Depends(get_db),
):
    try:
        payload = _require_token(token)
        user_id = payload.get("user_id")
        file = db.query(FileModel).filter(FileModel.id == file_id).first()

        if not file or file.user_id != user_id:
            raise HTTPException(status_code=403, detail="Нет доступа")

        if not request.name or request.name.strip() == "":
            raise HTTPException(
                status_code=400, detail="Имя файла не может быть пустым"
            )

        file.filename = request.name
        db.commit()
        db.refresh(file)
        return {"ok": True, "name": file.filename}
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@router.put("/rename-folder/{folder_id}")
def rename_folder(
    folder_id: str,
    request: RenameRequest,
    token: str = Query(...),
    db: Session = Depends(get_db),
):
    try:
        payload = _require_token(token)
        user_id = payload.get("user_id")

        folder = (
            db.query(FileModel)
            .filter(
                FileModel.user_id == user_id,
                FileModel.folder_id == folder_id,
                FileModel.file_type == "folder",
            )
            .first()
        )

        if not folder:
            raise HTTPException(status_code=403, detail="Нет доступа")

        if not request.name or request.name.strip() == "":
            raise HTTPException(
                status_code=400, detail="Имя папки не может быть пустым"
            )

        folder.filename = request.name
        db.commit()
        db.refresh(folder)
        return {"ok": True, "name": folder.filename}
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/upload")
async def upload_file(
    token: str = Query(...),
    folder_id: str = Form(...),
    is_public: bool = Form(False),
    file: UploadFile = FastAPIFile(...),
    db: Session = Depends(get_db),
):
    """Завантажити файл у папку"""
    try:
        payload = _require_token(token)
        user_id = payload.get("user_id")
        user = db.query(User).filter(User.id == user_id).first()

        if not user:
            raise HTTPException(status_code=401, detail="Пользователь не найден")

        if not file.filename:
            raise HTTPException(status_code=400, detail="Имя файла отсутствует")

        # Генерувати унікальне імʼя для файлу
        file_ext = Path(file.filename).suffix
        unique_filename = f"{uuid.uuid4()}{file_ext}"
        file_path = UPLOAD_DIR / unique_filename

        # Зберегти файл
        contents = await file.read()
        if not contents:
            raise HTTPException(
                status_code=400, detail="Пустой файл или ошибка загрузки"
            )
        with open(file_path, "wb") as f:
            f.write(contents)

        file_type = get_file_type(file.filename)
        file_size = len(contents)

        # Зберегти інформацію в БД
        db_file = FileModel(
            user_id=user.id,
            folder_id=folder_id,
            filename=file.filename,
            file_path=str(file_path),
            file_type=file_type,
            file_size=file_size,
            is_public=is_public,
            created_at=datetime.utcnow(),
        )
        db.add(db_file)
        db.commit()
        db.refresh(db_file)

        return {
            "id": db_file.id,
            "name": db_file.filename,
            "size": f"{db_file.file_size / 1024 / 1024:.1f} MB",
            "type": db_file.file_type,
            "created_at": db_file.created_at.isoformat(),
        }

    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        import traceback

        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Ошибка: {str(e)}")


@router.get("/folder/{folder_id}")
def get_folder_files(
    folder_id: str, token: str = Query(None), db: Session = Depends(get_db)
):
    """Отримати всі файли папки (публічні або приватні з токеном)"""
    try:
        user_id = None

        # Спробуємо отримати user_id з токена, якщо він є
        if token:
            try:
                payload = decode_access_token(token)
                user_id = payload.get("user_id")
            except ValueError:
                pass  # Токен невірний, але можемо показати публічні файли

        # Отримуємо саму папку
        folder = (
            db.query(FileModel)
            .filter(FileModel.folder_id == folder_id, FileModel.file_type == "folder")
            .first()
        )

        if not folder:
            raise HTTPException(status_code=404, detail="Папка не найдена")

        # Перевіряємо доступ до папки
        is_owner = user_id == folder.user_id
        is_public = folder.is_public

        if not is_owner and not is_public:
            raise HTTPException(status_code=403, detail="Нет доступа к папке")

        user = db.query(User).filter(User.id == folder.user_id).first()
        if not user:
            raise HTTPException(status_code=404, detail="Пользователь не найден")

        # Отримуємо файли з папки
        if is_owner:
            # Власник бачить всі файли
            files = (
                db.query(FileModel)
                .filter(
                    FileModel.user_id == user_id,
                    FileModel.folder_id == folder_id,
                    FileModel.file_type != "folder",
                )
                .all()
            )
        else:
            # Для публічної папки показуємо ВСІ файли (папка публічна = весь вміст видимий)
            files = (
                db.query(FileModel)
                .filter(
                    FileModel.folder_id == folder_id, FileModel.file_type != "folder"
                )
                .all()
            )

        # Визначаємо, чи може користувач видалити ПАПКУ:
        # Публічна папка: можуть видаляти ТІЛЬКИ адміністратори
        # Приватна папка: може видаляти ТІЛЬКИ власник
        can_delete_folder = False
        if token and user_id:
            current_user = db.query(User).filter(User.id == user_id).first()
            if current_user:
                if is_public:
                    # Публічну: тільки адміністратор
                    can_delete_folder = current_user.role == "admin"
                else:
                    # Приватну: тільки власник
                    can_delete_folder = is_owner

        return {
            "can_delete_folder": can_delete_folder,
            "folder_owner_id": folder.user_id,
            "files": [
                {
                    "id": f.id,
                    "name": f.filename,
                    "size": f"{f.file_size / 1024 / 1024:.1f} MB"
                    if f.file_type != "folder"
                    else "-",
                    "type": f.file_type,
                    "uploaded_at": f.created_at.isoformat(),
                    "is_public": f.is_public if f.file_type != "folder" else None,
                    "can_delete": is_owner or (token and user_id == f.user_id),
                }
                for f in files
            ],
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.delete("/delete-folder/{folder_id}")
def delete_folder(
    folder_id: str, token: str = Query(...), db: Session = Depends(get_db)
):
    """Видалити папку та всі файли в ній

    Дозволяє:
    - власнику видалити свою папку
    - адміністратору видалити публічну папку будь-кого
    """
    try:
        payload = _require_token(token)
        user_id = payload.get("user_id")

        # Отримаємо користувача для перевірки ролі
        current_user = db.query(User).filter(User.id == user_id).first()
        if not current_user:
            raise HTTPException(status_code=401, detail="Пользователь не найден")

        # Шукаємо папку без фільтра по user_id (адмін може видалити будь-яку публічну)
        folder_record = (
            db.query(FileModel)
            .filter(
                FileModel.folder_id == folder_id,
                FileModel.file_type == "folder",
            )
            .first()
        )

        if not folder_record:
            raise HTTPException(status_code=404, detail="Папка не найдена")

        # Перевіряємо права доступу:
        # Публічні папки: видяляє ТІЛЬКИ адміністратор
        # Приватні папки: видяляє ТІЛЬКИ власник
        is_owner = folder_record.user_id == user_id
        is_admin = current_user.role == "admin"

        if folder_record.is_public:
            # Публічну папку може видалити ТІЛЬКИ адміністратор
            if not is_admin:
                raise HTTPException(
                    status_code=403,
                    detail="Публичные папки могут удалять только администраторы",
                )
        else:
            # Приватну папку може видалити ТІЛЬКИ власник
            if not is_owner:
                raise HTTPException(
                    status_code=403,
                    detail="Вы не имеете доступа к удалению этой папки",
                )

        # Видаляємо всі файли у папці
        files = db.query(FileModel).filter(FileModel.folder_id == folder_id).all()

        for f in files:
            if f.file_path and os.path.exists(f.file_path) and f.file_type != "folder":
                os.remove(f.file_path)

        for f in files:
            db.delete(f)

        # Видаляємо саму папку з диску
        folder_path = Path(folder_record.file_path)
        if folder_path.exists() and folder_path.is_dir():
            folder_path.rmdir()

        db.commit()

        return {"ok": True, "message": "Папка удалена"}

    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@router.api_route("/download/{file_id}", methods=["GET", "HEAD"])
def download_file(
    file_id: int,
    token: str = Query(None),
    as_pdf: bool = Query(False),
    db: Session = Depends(get_db),
):
    try:
        user_id = None
        if token:
            try:
                payload = decode_access_token(token)
                user_id = payload.get("user_id")
            except ValueError:
                raise HTTPException(
                    status_code=401, detail="Неверный или просроченный токен"
                )

        file = db.query(FileModel).filter(FileModel.id == file_id).first()
        if not file:
            raise HTTPException(status_code=404, detail="Файл не найден")

        # Перевіряємо доступ: власник файлу або публічний файл
        is_owner = user_id and file.user_id == user_id
        is_public = getattr(file, "is_public", False)

        if not (is_owner or is_public):
            raise HTTPException(status_code=403, detail="Нет доступа к файлу")

        if not os.path.exists(file.file_path):
            raise HTTPException(status_code=404, detail="Файл на диске не найден")

        ext = Path(file.filename).suffix.lower()

        # If requested, convert PPT/PPTX to PDF for inline browser viewing
        if as_pdf and ext in {".pptx", ".ppt"}:
            pdf_path = Path(file.file_path).with_suffix(".pdf")
            try:
                if not pdf_path.exists():
                    convert_ppt_to_pdf(file.file_path, str(pdf_path))
                response = FileResponse(str(pdf_path), media_type="application/pdf")
                response.headers["Content-Disposition"] = (
                    f'inline; filename="{file.filename.rsplit(".", 1)[0]}.pdf"'
                )
                return response
            except Exception as e:
                raise HTTPException(status_code=500, detail=str(e))

        text_ext = {
            ".txt",
            ".log",
            ".ini",
            ".cfg",
            ".yaml",
            ".yml",
            ".csv",
            ".html",
            ".json",
            ".js",
            ".py",
            ".cpp",
            ".xml",
            ".css",
            ".md",
            ".rtf",
        }

        if ext in text_ext:
            media_type = "text/plain; charset=utf-8"
            disposition = "inline"

        elif ext == ".pdf":
            media_type = "application/pdf"
            disposition = "inline"

        elif ext == ".rtf":
            media_type = "application/rtf"
            disposition = "inline"

        elif ext == ".odt":
            media_type = "application/vnd.oasis.opendocument.text"
            disposition = "inline"

        elif ext == ".pptx":
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            disposition = "inline"

        elif ext == ".ppt":
            media_type = "application/vnd.ms-powerpoint"
            disposition = "inline"

        else:
            media_type_guess, _ = mimetypes.guess_type(file.filename)
            media_type = media_type_guess or "application/octet-stream"
            disposition = "attachment"

        # Build a safe Content-Disposition header that supports non-ASCII file names.
        # Starlette's default header builder can fail when we set the header manually
        # with a raw unicode value (e.g. Ukrainian/Cyrillic file names).
        # We use RFC 5987 encoding (filename*) for UTF-8 names.
        try:
            # If filename is ASCII-only, we can use the simple form.
            file.filename.encode("ascii")
            disposition_value = f'{disposition}; filename="{file.filename}"'
        except UnicodeEncodeError:
            disposition_value = (
                f"{disposition}; filename*=UTF-8''{quote(file.filename)}"
            )

        response = FileResponse(
            file.file_path,
            media_type=media_type,
        )
        response.headers["Content-Disposition"] = disposition_value
        return response

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/public")
def get_public_files(request: Request, db: Session = Depends(get_db)):
    try:
        user_id = None
        is_admin = False

        # Отримуємо токен з Authorization header
        auth_header = request.headers.get("Authorization", "")
        token = None
        if auth_header.startswith("Bearer "):
            token = auth_header[7:]  # Видаляємо "Bearer " префікс

        if token:
            try:
                payload = decode_access_token(token)
                user_id = payload.get("user_id")
                user = db.query(User).filter(User.id == user_id).first()
                if user and user.role == "admin":
                    is_admin = True
                print(
                    f"[DEBUG] Token received, user_id={user_id}, user={user}, is_admin={is_admin}"
                )
            except ValueError as e:
                print(f"[DEBUG] Token decode error: {e}")
                pass  # Invalid token, but we can still show public files
        else:
            print("[DEBUG] No token provided")

        if not hasattr(FileModel, "is_public"):
            return []
        # Отримуємо всі публічні об'єкти (файли та папки)
        items = db.query(FileModel).filter(FileModel.is_public == True).all()

        result = []
        for f in items:
            # Для папок: підраховуємо файли, середню оцінку, кількість оцінок і коментарів
            if f.file_type == "folder":
                file_count = (
                    db.query(FileModel)
                    .filter(
                        FileModel.folder_id == f.folder_id,
                        FileModel.file_type != "folder",
                    )
                    .count()
                )

                # Середня оцінка та кількість оцінок
                rating_result = (
                    db.query(func.avg(Rating.rating), func.count(Rating.id))
                    .filter(Rating.folder_id == f.folder_id)
                    .first()
                )
                average_rating = float(rating_result[0]) if rating_result[0] else 0.0
                rating_count = rating_result[1] or 0

                # Кількість коментарів
                comment_count = (
                    db.query(Comment).filter(Comment.folder_id == f.folder_id).count()
                )

                can_delete = is_admin  # Only admin can delete public folders
            else:
                # Для файлів: ці поля не застосовні, встановлюємо 0
                file_count = 0
                average_rating = 0.0
                rating_count = 0
                comment_count = 0
                can_delete = False  # Files are not deleted from public page

            result.append(
                {
                    "id": f.id if f.file_type != "folder" else f.folder_id,
                    "name": f.filename,
                    "size": f"{f.file_size / 1024 / 1024:.1f} MB"
                    if f.file_type != "folder"
                    else "-",
                    "type": f.file_type,
                    "uploaded_at": f.created_at.isoformat(),
                    "user_id": f.user_id,
                    "is_public": f.is_public,
                    "folder_id": f.folder_id,
                    "file_count": file_count,
                    "average_rating": average_rating,
                    "rating_count": rating_count,
                    "comment_count": comment_count,
                    "can_delete": can_delete,
                }
            )

        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/folder/{folder_id}/set-public")
def set_folder_public(
    folder_id: str,
    is_public: bool = Form(...),
    token: str = Query(...),
    db: Session = Depends(get_db),
):
    try:
        payload = _require_token(token)
        user_id = payload.get("user_id")

        folder = (
            db.query(FileModel)
            .filter(
                FileModel.user_id == user_id,
                FileModel.folder_id == folder_id,
                FileModel.file_type == "folder",
            )
            .first()
        )

        if not folder:
            raise HTTPException(status_code=403, detail="Нет доступа")

        folder.is_public = bool(is_public)
        db.commit()
        return {"ok": True, "is_public": folder.is_public}
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/file/{file_id}/set-public")
def set_file_public(
    file_id: int,
    is_public: bool = Form(...),
    token: str = Query(...),
    db: Session = Depends(get_db),
):
    try:
        payload = _require_token(token)
        user_id = payload.get("user_id")
        file = db.query(FileModel).filter(FileModel.id == file_id).first()
        if not file or file.user_id != user_id:
            raise HTTPException(status_code=403, detail="Нет доступа")
        if not hasattr(file, "is_public"):
            raise HTTPException(
                status_code=400, detail="Поле is_public не настроено в модели"
            )
        file.is_public = bool(is_public)
        db.commit()
        return {"ok": True, "is_public": file.is_public}
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


# =============== ПЕРЕМІЩЕННЯ ФАЙЛІВ ===============
@router.post("/move-file/{file_id}")
def move_file(
    file_id: int,
    target_folder_id: str = Form(...),
    token: str = Query(...),
    db: Session = Depends(get_db),
):
    """Переместить файл в іншу папку"""
    try:
        payload = _require_token(token)
        user_id = payload.get("user_id")

        # Перевіряємо, що файл існує і належить користувачу
        file = db.query(FileModel).filter(FileModel.id == file_id).first()
        if not file or file.user_id != user_id:
            raise HTTPException(status_code=403, detail="Нет доступа к файлу")

        # Перевіряємо, що цільова папка існує і належить користувачу
        target_folder = (
            db.query(FileModel)
            .filter(
                FileModel.user_id == user_id,
                FileModel.folder_id == target_folder_id,
                FileModel.file_type == "folder",
            )
            .first()
        )

        if not target_folder:
            raise HTTPException(
                status_code=403, detail="Целевая папка не найдена или нет доступа"
            )

        # Не можна переміщувати файл в тю ж п папку
        if file.folder_id == target_folder_id:
            raise HTTPException(status_code=400, detail="Файл уже в этой папке")

        # Переміщуємо файл
        file.folder_id = target_folder_id
        db.commit()
        db.refresh(file)

        return {"ok": True, "message": "Файл успешно перемещен"}

    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


# =============== ВИДАЛЕННЯ ФАЙЛІВ ===============
@router.delete("/file/{file_id}")
def delete_file(file_id: int, token: str = Query(...), db: Session = Depends(get_db)):
    """Видалити файл (тільки власник)

    Повертає ok: true, якщо видалено.
    """
    try:
        payload = _require_token(token)
        user_id = payload.get("user_id")

        file = db.query(FileModel).filter(FileModel.id == file_id).first()
        if not file:
            raise HTTPException(status_code=404, detail="Файл не найден")

        # Дозволено лише власнику файлу
        if file.user_id != user_id:
            raise HTTPException(status_code=403, detail="Нет доступа к файлу")

        # Видаляємо файл з диску
        if file.file_path and os.path.exists(file.file_path):
            try:
                os.remove(file.file_path)
            except Exception:
                pass

        db.delete(file)
        db.commit()

        return {"ok": True, "message": "Файл удален"}
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


# =============== ЧИТАННЯ ЗМІСТУ ДОКУМЕНТІВ ===============
@router.get("/read/{file_id}")
def read_file_content(
    file_id: int, token: str = Query(None), db: Session = Depends(get_db)
):
    """Прочитати вміст документу (DOCX, DOC, TXT, JSON, JS) та повернути як JSON"""
    try:
        user_id = None
        if token:
            try:
                payload = decode_access_token(token)
                user_id = payload.get("user_id")
            except ValueError:
                raise HTTPException(
                    status_code=401, detail="Неверный или просроченный токен"
                )

        file = db.query(FileModel).filter(FileModel.id == file_id).first()
        if not file:
            raise HTTPException(status_code=404, detail="Файл не найден")

        # Перевіряємо доступ: власник файлу або публічний файл
        is_owner = user_id and file.user_id == user_id
        is_public = getattr(file, "is_public", False)

        if not (is_owner or is_public):
            raise HTTPException(status_code=403, detail="Нет доступа к файлу")

        if not os.path.exists(file.file_path):
            raise HTTPException(status_code=404, detail="Файл на диске не найден")

        file_ext = Path(file.filename).suffix.lower()
        content = ""

        # Читання DOCX файлів
        if file_ext == ".docx":
            try:
                doc = DocxDocument(file.file_path)
                paragraphs = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        paragraphs.append(para.text)
                content = "\n".join(paragraphs)
            except Exception as e:
                raise HTTPException(
                    status_code=400, detail=f"Ошибка чтения DOCX: {str(e)}"
                )

        # Читання ODT файлів
        elif file_ext == ".odt":
            try:
                with zipfile.ZipFile(file.file_path, "r") as zf:
                    with zf.open("content.xml") as content_xml:
                        tree = ET.parse(content_xml)
                        root = tree.getroot()
                        ns = {"text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0"}
                        paragraphs = []
                        for elem in root.findall(".//text:p", ns):
                            if elem.text and elem.text.strip():
                                paragraphs.append(elem.text.strip())
                        content = "\n".join(paragraphs)
            except Exception as e:
                raise HTTPException(
                    status_code=400, detail=f"Ошибка чтения ODT: {str(e)}"
                )

        # Читання PPTX файлів
        elif file_ext == ".pptx":
            try:
                prs = Presentation(file.file_path)
                paragraphs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            paragraphs.append(shape.text.strip())
                content = "\n".join([p for p in paragraphs if p])
            except Exception as e:
                raise HTTPException(
                    status_code=400, detail=f"Ошибка чтения PPTX: {str(e)}"
                )

        # Читання PPT (старих .ppt) файлів (грубий текст)
        elif file_ext == ".ppt":
            try:
                with open(file.file_path, "rb") as f:
                    content_bytes = f.read()
                content = "".join(
                    chr(b) for b in content_bytes if 32 <= b < 127 or b in {9, 10, 13}
                )
                content = "\n".join(
                    line.strip() for line in content.split("\n") if line.strip()
                )
            except Exception as e:
                raise HTTPException(
                    status_code=400, detail=f"Нельзя прочитать PPT файл: {str(e)}"
                )

        # Читання інших текстових форматів
        elif file_ext in {
            ".txt",
            ".log",
            ".ini",
            ".cfg",
            ".yaml",
            ".yml",
            ".csv",
            ".json",
            ".js",
            ".py",
            ".html",
            ".xml",
            ".css",
            ".md",
            ".cpp",
            ".rtf",
        }:
            try:
                with open(file.file_path, "r", encoding="utf-8") as f:
                    content = f.read()
            except UnicodeDecodeError:
                try:
                    with open(file.file_path, "r", encoding="cp1251") as f:
                        content = f.read()
                except UnicodeDecodeError:
                    try:
                        with open(file.file_path, "r", encoding="latin-1") as f:
                            content = f.read()
                    except Exception as e:
                        raise HTTPException(
                            status_code=400,
                            detail=f"Не удалось прочитать файл в известных кодировках (utf-8, cp1251, latin-1): {str(e)}",
                        )
            except Exception as e:
                raise HTTPException(
                    status_code=400, detail=f"Ошибка чтения файла: {str(e)}"
                )

        # Для DOC (старих Word файлів) - спробуємо витягти текст
        elif file_ext == ".doc":
            # DOC формат був замінений на DOCX,
            # спробуємо прочитати як текст із кодуванням
            try:
                with open(file.file_path, "rb") as f:
                    content_bytes = f.read()
                # Спробуємо витягти ASCII символи
                content = "".join(
                    chr(b) for b in content_bytes if 32 <= b < 127 or b in {9, 10, 13}
                )
                # Очистимо білі символи
                content = "\n".join(
                    line.strip() for line in content.split("\n") if line.strip()
                )
            except Exception as e:
                raise HTTPException(
                    status_code=400, detail=f"Нельзя прочитать DOC файл: {str(e)}"
                )

        else:
            raise HTTPException(
                status_code=400,
                detail=f"Формат {file_ext} не поддерживается. Поддерживаются: .docx, .doc, .odt, .pptx, .ppt, .txt, .json, .js, .py, .html та інші текстові формати",
            )

        return {
            "file_name": file.filename,
            "file_type": file.file_type,
            "file_extension": file_ext,
            "file_size_bytes": file.file_size,
            "content": content,
            "content_length": len(content),
            "created_at": file.created_at.isoformat(),
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/preview/{file_id}", response_class=HTMLResponse)
def preview_file(file_id: int, token: str = Query(None), db: Session = Depends(get_db)):
    """
    Превью текстового содержимого (DOCX, TXT и др.) в виде HTML.
    Использует уже существную функцию read_file_content для извлечения текста.
    """
    result = read_file_content(file_id, token, db)
    content = result.get("content", "")
    escaped = pyhtml.escape(content).replace("\n", "<br/>")
    html_page = f"<html><head><meta charset='utf-8'></head><body style='font-family:Arial, sans-serif;white-space:pre-wrap'>{escaped}</body></html>"
    return HTMLResponse(content=html_page, status_code=200)
    html_page = f"<html><head><meta charset='utf-8'></head><body style='font-family:Arial, sans-serif;white-space:pre-wrap'>{escaped}</body></html>"
    return HTMLResponse(content=html_page, status_code=200)
    html_page = f"<html><head><meta charset='utf-8'></head><body style='font-family:Arial, sans-serif;white-space:pre-wrap'>{escaped}</body></html>"
    return HTMLResponse(content=html_page, status_code=200)
    html_page = f"<html><head><meta charset='utf-8'></head><body style='font-family:Arial, sans-serif;white-space:pre-wrap'>{escaped}</body></html>"
    return HTMLResponse(content=html_page, status_code=200)
    html_page = f"<html><head><meta charset='utf-8'></head><body style='font-family:Arial, sans-serif;white-space:pre-wrap'>{escaped}</body></html>"
    return HTMLResponse(content=html_page, status_code=200)
    html_page = f"<html><head><meta charset='utf-8'></head><body style='font-family:Arial, sans-serif;white-space:pre-wrap'>{escaped}</body></html>"
    return HTMLResponse(content=html_page, status_code=200)
    html_page = f"<html><head><meta charset='utf-8'></head><body style='font-family:Arial, sans-serif;white-space:pre-wrap'>{escaped}</body></html>"
    return HTMLResponse(content=html_page, status_code=200)
    html_page = f"<html><head><meta charset='utf-8'></head><body style='font-family:Arial, sans-serif;white-space:pre-wrap'>{escaped}</body></html>"
    return HTMLResponse(content=html_page, status_code=200)
    html_page = f"<html><head><meta charset='utf-8'></head><body style='font-family:Arial, sans-serif;white-space:pre-wrap'>{escaped}</body></html>"
    return HTMLResponse(content=html_page, status_code=200)
    html_page = f"<html><head><meta charset='utf-8'></head><body style='font-family:Arial, sans-serif;white-space:pre-wrap'>{escaped}</body></html>"
    return HTMLResponse(content=html_page, status_code=200)
